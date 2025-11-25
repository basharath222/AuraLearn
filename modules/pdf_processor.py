import pdfplumber
import PyPDF2
import pandas as pd
import docx
from pptx import Presentation
import io

def extract_text_from_pdf(uploaded_file):
    """
    Extracts text based on file extension (PDF, DOCX, PPTX, CSV, XLSX, XML, TXT).
    Uses file extensions instead of MIME types for better reliability.
    """
    text = ""
    # Convert filename to lowercase to handle .PDF or .pdf
    filename = uploaded_file.name.lower()

    try:
        # 1. Handle PDF
        if filename.endswith(".pdf"):
            # Try with pdfplumber first
            try:
                with pdfplumber.open(uploaded_file) as pdf:
                    for page in pdf.pages:
                        content = page.extract_text()
                        if content:
                            text += content + "\n"
            except Exception as e:
                print(f"pdfplumber failed: {e}")

            # Fallback to PyPDF2 if text seems too short or empty
            if len(text.strip()) < 200:
                try:
                    uploaded_file.seek(0)  # Reset file pointer
                    reader = PyPDF2.PdfReader(uploaded_file)
                    for page in reader.pages:
                        content = page.extract_text()
                        if content:
                            text += content + "\n"
                except Exception as e:
                    return f"⚠️ Error reading PDF with fallback: {e}"
            
            return text.strip() if text.strip() else "⚠️ PDF scanned or empty."

        # 2. Handle Word (.docx)
        elif filename.endswith(".docx"):
            doc = docx.Document(uploaded_file)
            fullText = []
            for para in doc.paragraphs:
                fullText.append(para.text)
            return '\n'.join(fullText)

        # 3. Handle PowerPoint (.pptx)
        elif filename.endswith(".pptx"):
            prs = Presentation(uploaded_file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        # 4. Handle Text / Markdown (.txt, .md)
        elif filename.endswith(".txt") or filename.endswith(".md"):
            return uploaded_file.read().decode("utf-8")

        # 5. Handle CSV (.csv)
        elif filename.endswith(".csv"):
            uploaded_file.seek(0)
            df = pd.read_csv(uploaded_file)
            return df.to_string()

        # 6. Handle Excel (.xlsx, .xls)
        elif filename.endswith(".xlsx") or filename.endswith(".xls"):
            uploaded_file.seek(0)
            df = pd.read_excel(uploaded_file)
            return df.to_string()
        
        # 7. Handle XML (.xml)
        elif filename.endswith(".xml"):
            return uploaded_file.read().decode("utf-8")

        else:
            return f"⚠️ Unsupported file type: {filename}. Please upload PDF, DOCX, PPTX, TXT, CSV, or Excel."

    except Exception as e:
        return f"⚠️ Error processing file: {str(e)}"